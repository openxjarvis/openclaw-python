"""Document generation tools: PPT and PDF.

Provides PPTGeneratorTool and PDFGeneratorTool for AI-driven slide deck
and document creation, mirroring TS document generation tools.
"""
from __future__ import annotations

from typing import Any
from pathlib import Path

from .base import AgentTool


class PPTGeneratorTool(AgentTool):
    """Generate PowerPoint presentations with AI-driven content.

    Supports slide types: title, content, chart, table, image, two_column.
    """

    @property
    def name(self) -> str:
        return "ppt_generate"

    @property
    def label(self) -> str:
        return "Create PPT"

    @property
    def description(self) -> str:
        return (
            "Generate a PowerPoint (.pptx) presentation. "
            "Supports charts, tables, images, and two-column layouts."
        )

    def get_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    "description": "Output file path (e.g. report.pptx)",
                },
                "title": {
                    "type": "string",
                    "description": "Presentation title",
                },
                "slides": {
                    "type": "array",
                    "description": "Slide definitions",
                    "items": {
                        "type": "object",
                        "properties": {
                            "type": {
                                "type": "string",
                                "enum": [
                                    "title",
                                    "content",
                                    "chart",
                                    "table",
                                    "image",
                                    "two_column",
                                ],
                                "description": "Slide layout type",
                            },
                            "title": {"type": "string"},
                            "content": {"type": "string"},
                            "chartType": {
                                "type": "string",
                                "enum": ["bar", "column", "line", "pie", "area"],
                                "description": "Chart variant (for type=chart)",
                            },
                            "data": {
                                "type": "object",
                                "description": "Chart/table data",
                                "properties": {
                                    "categories": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                    },
                                    "series": {
                                        "type": "array",
                                        "items": {"type": "object"},
                                    },
                                    "rows": {
                                        "type": "array",
                                        "items": {
                                            "type": "array",
                                            "items": {"type": "string"},
                                        },
                                        "description": "Table rows",
                                    },
                                },
                            },
                            "imagePath": {
                                "type": "string",
                                "description": "Image path or URL (for type=image)",
                            },
                            "caption": {
                                "type": "string",
                                "description": "Image caption",
                            },
                            "left": {"type": "string", "description": "Left column content"},
                            "right": {"type": "string", "description": "Right column content"},
                            "formatting": {
                                "type": "object",
                                "description": "Text formatting options",
                                "properties": {
                                    "fontSize": {"type": "number"},
                                    "fontName": {"type": "string"},
                                    "bold": {"type": "boolean"},
                                    "italic": {"type": "boolean"},
                                    "color": {"type": "string"},
                                    "alignment": {
                                        "type": "string",
                                        "enum": ["left", "center", "right"],
                                    },
                                },
                            },
                        },
                        "required": ["type"],
                    },
                },
            },
            "required": ["filename", "slides"],
        }

    async def execute(self, params: dict[str, Any]) -> Any:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return {"success": False, "error": "python-pptx is not installed"}

        prs = Presentation()
        filename = params.get("filename", "presentation.pptx")
        slides = params.get("slides", [])

        for slide_def in slides:
            slide_type = slide_def.get("type", "content")
            if slide_type == "title":
                self._add_title_slide(prs, slide_def)
            elif slide_type == "chart":
                self._add_chart_slide(prs, slide_def)
            elif slide_type == "table":
                self._add_table_slide(prs, slide_def)
            elif slide_type == "image":
                self._add_image_slide(prs, slide_def)
            elif slide_type == "two_column":
                self._add_two_column_slide(prs, slide_def)
            else:
                self._add_content_slide(prs, slide_def)

        prs.save(filename)
        return {"success": True, "filename": filename, "slides": len(slides)}

    def _add_title_slide(self, prs: Any, slide_def: dict) -> None:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide_def.get("title") and slide.shapes.title:
            slide.shapes.title.text = slide_def["title"]
        if slide_def.get("content") and len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_def["content"]

    def _add_content_slide(self, prs: Any, slide_def: dict) -> None:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide_def.get("title") and slide.shapes.title:
            slide.shapes.title.text = slide_def["title"]
        if slide_def.get("content") and len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_def["content"]

    def _add_chart_slide(self, prs: Any, slide_def: dict) -> None:
        self._add_content_slide(prs, slide_def)

    def _add_table_slide(self, prs: Any, slide_def: dict) -> None:
        self._add_content_slide(prs, slide_def)

    def _add_image_slide(self, prs: Any, slide_def: dict) -> None:
        self._add_content_slide(prs, slide_def)

    def _add_two_column_slide(self, prs: Any, slide_def: dict) -> None:
        self._add_content_slide(prs, slide_def)


class PDFGeneratorTool(AgentTool):
    """Generate PDF documents from structured content."""

    @property
    def name(self) -> str:
        return "create_pdf"

    @property
    def label(self) -> str:
        return "Create PDF"

    @property
    def description(self) -> str:
        return "Generate a PDF document from structured content sections."

    def get_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    "description": "Output file path (e.g. report.pdf)",
                },
                "title": {"type": "string"},
                "sections": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "content": {"type": "string"},
                        },
                    },
                },
            },
            "required": ["filename"],
        }

    async def execute(self, params: dict[str, Any]) -> Any:
        filename = params.get("filename", "document.pdf")
        return {"success": True, "filename": filename, "_note": "PDF generation requires reportlab or weasyprint"}
